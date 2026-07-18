import os
import zipfile


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".png", ".jpg", ".jpeg"}

MAX_FILE_BYTES = 100 * 1024 * 1024
MAX_OFFICE_FILE_BYTES = 25 * 1024 * 1024
MAX_ARCHIVE_ENTRIES = 10_000
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 100 * 1024 * 1024
MAX_ARCHIVE_ENTRY_BYTES = 25 * 1024 * 1024
MAX_ARCHIVE_RATIO = 200
MAX_CHARS = 20_000

MAX_PDF_PAGES = 10
MAX_DOCX_PARAGRAPHS = 1_000
MAX_DOCX_TABLES = 50
MAX_DOCX_TABLE_ROWS = 200
MAX_XLSX_SHEETS = 8
MAX_XLSX_ROWS_PER_SHEET = 80
MAX_XLSX_SCANNED_ROWS = 500
MAX_XLSX_SCANNED_COLUMNS = 100
MAX_PPTX_SLIDES = 20
MAX_PPTX_SHAPES_PER_SLIDE = 250


class DocumentPreviewError(Exception):
    pass


def _validate_file(file_path):
    if not isinstance(file_path, str) or not file_path:
        raise DocumentPreviewError("Missing file path")
    if not os.path.isfile(file_path):
        raise DocumentPreviewError("File not found")

    extension = os.path.splitext(file_path)[1].lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise DocumentPreviewError("Unsupported document type")

    file_size = os.path.getsize(file_path)
    max_size = MAX_OFFICE_FILE_BYTES if extension != ".pdf" else MAX_FILE_BYTES
    if file_size > max_size:
        raise DocumentPreviewError(
            f"Document is too large to preview safely ({file_size // (1024 * 1024)} MB)"
        )

    if extension in {".docx", ".xlsx", ".pptx"}:
        _validate_office_archive(file_path)

    return extension


def _validate_office_archive(file_path):
    if not zipfile.is_zipfile(file_path):
        raise DocumentPreviewError("Office document is invalid or corrupted")

    total_uncompressed = 0
    try:
        with zipfile.ZipFile(file_path) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_ARCHIVE_ENTRIES:
                raise DocumentPreviewError("Office document contains too many archive entries")

            for entry in entries:
                if entry.file_size > MAX_ARCHIVE_ENTRY_BYTES:
                    raise DocumentPreviewError("Office document contains an oversized embedded item")
                total_uncompressed += entry.file_size
                if total_uncompressed > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
                    raise DocumentPreviewError("Office document expands beyond the safe preview limit")
                if entry.compress_size and entry.file_size / entry.compress_size > MAX_ARCHIVE_RATIO:
                    raise DocumentPreviewError("Office document has an unsafe compression ratio")
    except zipfile.BadZipFile as exc:
        raise DocumentPreviewError("Office document is invalid or corrupted") from exc


def _finalize(parts, kind, truncated=False):
    content = "\n\n".join(part for part in parts if part)
    was_truncated = truncated or len(content) > MAX_CHARS
    return {
        "success": True,
        "content": content[:MAX_CHARS],
        "kind": kind,
        "truncated": was_truncated,
    }


def _preview_pdf(file_path):
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    if getattr(reader, "is_encrypted", False):
        raise DocumentPreviewError("This PDF is password protected")

    parts = []
    truncated = len(reader.pages) > MAX_PDF_PAGES
    for page_number, page in enumerate(reader.pages, start=1):
        if page_number > MAX_PDF_PAGES:
            break
        page_text = (page.extract_text() or "").strip()
        try:
            has_images = bool(page.images)
        except Exception:
            has_images = False

        if page_text:
            part = f"Page {page_number}\n{page_text}"
            if has_images:
                part += "\n\n[Image / graphic present on this page]"
        elif has_images:
            part = f"Page {page_number}\n[Scanned or image-based content: text could not be extracted]"
        else:
            part = f"Page {page_number}\n[No extractable text detected on this page]"
        parts.append(part)
        if sum(len(item) for item in parts) >= MAX_CHARS:
            truncated = True
            break
    return _finalize(parts, "PDF", truncated)


def _preview_docx(file_path):
    from docx import Document

    document = Document(file_path)
    parts = []
    truncated = len(document.paragraphs) > MAX_DOCX_PARAGRAPHS or len(document.tables) > MAX_DOCX_TABLES

    for paragraph in document.paragraphs[:MAX_DOCX_PARAGRAPHS]:
        paragraph_text = paragraph.text.strip()
        has_image = "w:drawing" in paragraph._p.xml or "w:pict" in paragraph._p.xml
        if paragraph_text:
            parts.append(paragraph_text)
        if has_image:
            parts.append("[Image / graphic]")

    for table in document.tables[:MAX_DOCX_TABLES]:
        if len(table.rows) > MAX_DOCX_TABLE_ROWS:
            truncated = True
        for row in table.rows[:MAX_DOCX_TABLE_ROWS]:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                has_image = "w:drawing" in cell._tc.xml or "w:pict" in cell._tc.xml
                if cell_text:
                    cells.append(cell_text)
                if has_image:
                    cells.append("[Image / graphic]")
            if cells:
                parts.append(" | ".join(cells))
    return _finalize(parts, "Word document", truncated)


def _preview_xlsx(file_path):
    from openpyxl import load_workbook

    workbook = load_workbook(file_path, read_only=False, data_only=False, keep_links=False)
    parts = []
    truncated = len(workbook.worksheets) > MAX_XLSX_SHEETS
    try:
        for sheet in workbook.worksheets[:MAX_XLSX_SHEETS]:
            sheet_parts = [f"Sheet: {sheet.title}"]
            populated_rows = 0
            max_row = min(sheet.max_row or 1, MAX_XLSX_SCANNED_ROWS)
            max_column = min(sheet.max_column or 1, MAX_XLSX_SCANNED_COLUMNS)
            if (sheet.max_row or 1) > max_row or (sheet.max_column or 1) > max_column:
                truncated = True

            for row in sheet.iter_rows(max_row=max_row, max_col=max_column):
                cells = [f"{cell.coordinate}: {cell.value}" for cell in row if cell.value is not None]
                if cells:
                    sheet_parts.append(" | ".join(cells))
                    populated_rows += 1
                if populated_rows >= MAX_XLSX_ROWS_PER_SHEET:
                    sheet_parts.append("[Additional rows not shown]")
                    truncated = True
                    break

            merged_ranges = list(sheet.merged_cells.ranges)
            if merged_ranges:
                shown_ranges = merged_ranges[:100]
                merged = ", ".join(str(cell_range) for cell_range in shown_ranges)
                sheet_parts.append(f"[Merged cells: {merged}]")
                if len(merged_ranges) > len(shown_ranges):
                    truncated = True
            if getattr(sheet, "_images", None):
                sheet_parts.append(f"[Images / graphics: {len(sheet._images)}]")
            if getattr(sheet, "_charts", None):
                sheet_parts.append(f"[Charts: {len(sheet._charts)}]")
            if populated_rows == 0 and not getattr(sheet, "_images", None) and not getattr(sheet, "_charts", None):
                sheet_parts.append("[No populated cells detected]")

            parts.append("\n".join(sheet_parts))
            if sum(len(item) for item in parts) >= MAX_CHARS:
                truncated = True
                break
    finally:
        workbook.close()
    return _finalize(parts, "Excel workbook", truncated)


def _preview_pptx(file_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(file_path)
    parts = []
    truncated = len(presentation.slides) > MAX_PPTX_SLIDES
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide_number > MAX_PPTX_SLIDES:
            break
        slide_parts = [f"Slide {slide_number}"]
        shapes = list(slide.shapes)
        if len(shapes) > MAX_PPTX_SHAPES_PER_SLIDE:
            truncated = True
        for shape in shapes[:MAX_PPTX_SHAPES_PER_SLIDE]:
            if getattr(shape, "has_text_frame", False):
                shape_text = shape.text.strip()
                if shape_text:
                    slide_parts.append(shape_text)

            if getattr(shape, "has_table", False):
                slide_parts.append("[Table]")
                for row_number, row in enumerate(shape.table.rows, start=1):
                    if row_number > MAX_DOCX_TABLE_ROWS:
                        truncated = True
                        break
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_parts.append("[Image / picture]")
            elif getattr(shape, "has_chart", False):
                slide_parts.append("[Chart]")
            elif shape.shape_type in (MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.MEDIA):
                slide_parts.append("[Graphic / media]")

        if len(slide_parts) == 1:
            slide_parts.append("[No extractable slide content detected]")
        parts.append("\n".join(slide_parts))
        if sum(len(item) for item in parts) >= MAX_CHARS:
            truncated = True
            break
    return _finalize(parts, "PowerPoint presentation", truncated)

def _preview_image(file_path):
    from core.extractor import extract_text
    text = extract_text(file_path)
    if text.strip():
        return _finalize([text], "Image (OCR)")
    else:
        return _finalize(["[No text detected in this image]"], "Image")


def build_document_preview(file_path):
    extension = _validate_file(file_path)
    try:
        if extension == ".pdf":
            return _preview_pdf(file_path)
        if extension == ".docx":
            return _preview_docx(file_path)
        if extension == ".xlsx":
            return _preview_xlsx(file_path)
        if extension in {".png", ".jpg", ".jpeg"}:
            return _preview_image(file_path)
        return _preview_pptx(file_path)
    except DocumentPreviewError:
        raise
    except Exception as exc:
        raise DocumentPreviewError(f"Could not preview document: {exc}") from exc
