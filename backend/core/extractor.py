from pypdf import PdfReader
from docx import Document

from core.access_control import check_read_access, classify_access_error

# Cap extracted text to ~100K chars (~20K words) to avoid huge files
# dominating indexing time. The chunker will further limit chunks.
_MAX_TEXT_CHARS = 100_000


def _extract_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" ".join(cells))
    wb.close()
    return "\n".join(parts)


def _extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" ".join(cells))
    return "\n".join(parts)

def _extract_image(path):
    try:
        from PIL import Image
        # Disable DecompressionBombWarning for large screenshots/scans
        Image.MAX_IMAGE_PIXELS = None
        
        from PIL.ExifTags import TAGS
        import winocr
        with Image.open(path) as img:
            # Check EXIF metadata for camera signatures
            try:
                exif_data = img._getexif()
                if exif_data:
                    for tag_id, value in exif_data.items():
                        tag = TAGS.get(tag_id, tag_id)
                        if tag in ('Make', 'Model', 'LensModel', 'Software'):
                            # Likely a camera photo or photo edited in Lightroom/Photoshop
                            # Skip OCR extraction
                            return ""
            except Exception:
                pass
            result = winocr.recognize_pil_sync(img)
            return result.get("text", "")
    except Exception as exc:
        import sys
        sys.stderr.write(f"[extractor] OCR failed for {path}: {exc}\n")
        sys.stderr.flush()
        return ""


def extract_text(path):
    try:
        lower_path = path.lower()
        if lower_path.endswith(".pdf"):
            reader = PdfReader(path)
            text = ""
            for page in reader.pages:
                text += (page.extract_text() or "") + " "
                if len(text) >= _MAX_TEXT_CHARS:
                    break
            return text[:_MAX_TEXT_CHARS]
        elif lower_path.endswith(".docx"):
            doc = Document(path)
            text = " ".join(paragraph.text for paragraph in doc.paragraphs)
            return text[:_MAX_TEXT_CHARS]
        elif lower_path.endswith(".xlsx") or lower_path.endswith(".xls"):
            return _extract_xlsx(path)[:_MAX_TEXT_CHARS]
        elif lower_path.endswith(".csv"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(_MAX_TEXT_CHARS)
        elif lower_path.endswith(".pptx"):
            return _extract_pptx(path)[:_MAX_TEXT_CHARS]
        elif lower_path.endswith((".png", ".jpg", ".jpeg")):
            return _extract_image(path)[:_MAX_TEXT_CHARS]        
        else:
            # Fallback for .txt, .md, .rtf, etc.
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(_MAX_TEXT_CHARS)
    except Exception:
        return ""            


def extract_text_with_status(path, allow_protected=False):
    if not allow_protected:
        readable, reason = check_read_access(path)
        if not readable:
            return "", reason

    try:
        lower_path = path.lower()
        if lower_path.endswith(".pdf"):
            reader = PdfReader(path)
            if getattr(reader, "is_encrypted", False):
                if not allow_protected:
                    return "", "password_protected"
                try:
                    reader.decrypt("")
                except Exception:
                    pass
            text = ""
            for page in reader.pages:
                text += (page.extract_text() or "") + " "
                if len(text) >= _MAX_TEXT_CHARS:
                    break
            return text[:_MAX_TEXT_CHARS], None
        elif lower_path.endswith(".docx"):
            doc = Document(path)
            text = " ".join(paragraph.text for paragraph in doc.paragraphs)
            return text[:_MAX_TEXT_CHARS], None
        elif lower_path.endswith(".xlsx") or lower_path.endswith(".xls"):
            return _extract_xlsx(path)[:_MAX_TEXT_CHARS], None
        elif lower_path.endswith(".csv"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(_MAX_TEXT_CHARS), None
        elif lower_path.endswith(".pptx"):
            return _extract_pptx(path)[:_MAX_TEXT_CHARS], None
        elif lower_path.endswith((".png", ".jpg", ".jpeg")):
            return _extract_image(path)[:_MAX_TEXT_CHARS], None
        else:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(_MAX_TEXT_CHARS), None
    except Exception as exc:
        reason = classify_access_error(exc)
        if reason:
            return "", reason
        message = str(exc).lower()
        if "password" in message or "encrypted" in message:
            return "", "password_protected"
        return "", "extract_error"