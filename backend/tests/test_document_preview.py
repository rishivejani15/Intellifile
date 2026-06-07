import os
import sys
import tempfile
import unittest
import zipfile
from unittest.mock import patch


BACKEND_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if BACKEND_DIR not in sys.path:
    sys.path.insert(0, BACKEND_DIR)

from core.document_preview import DocumentPreviewError, build_document_preview


class DocumentPreviewTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()

    def tearDown(self):
        self.temp_dir.cleanup()

    def path(self, name):
        return os.path.join(self.temp_dir.name, name)

    def test_pdf_preview_marks_page_without_extractable_text(self):
        from pypdf import PdfWriter

        file_path = self.path("blank.pdf")
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with open(file_path, "wb") as output:
            writer.write(output)

        result = build_document_preview(file_path)

        self.assertTrue(result["success"])
        self.assertEqual(result["kind"], "PDF")
        self.assertIn("[No extractable text detected on this page]", result["content"])

    def test_docx_preview_includes_paragraphs_and_tables(self):
        from docx import Document

        file_path = self.path("sample.docx")
        document = Document()
        document.add_paragraph("Quarterly report")
        table = document.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "Revenue"
        table.cell(0, 1).text = "42"
        document.save(file_path)

        result = build_document_preview(file_path)

        self.assertEqual(result["kind"], "Word document")
        self.assertIn("Quarterly report", result["content"])
        self.assertIn("Revenue | 42", result["content"])

    def test_xlsx_preview_is_bounded_for_large_dimensions(self):
        from openpyxl import Workbook

        file_path = self.path("large-dimensions.xlsx")
        workbook = Workbook()
        sheet = workbook.active
        sheet["A1"] = "Visible"
        sheet["A1000000"] = "Outside preview bounds"
        workbook.save(file_path)

        result = build_document_preview(file_path)

        self.assertEqual(result["kind"], "Excel workbook")
        self.assertIn("A1: Visible", result["content"])
        self.assertNotIn("Outside preview bounds", result["content"])
        self.assertTrue(result["truncated"])

    def test_pptx_preview_includes_text_and_table(self):
        from pptx import Presentation
        from pptx.util import Inches

        file_path = self.path("sample.pptx")
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Launch plan"
        table = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(5), Inches(1)).table
        table.cell(0, 0).text = "Owner"
        table.cell(0, 1).text = "Team"
        presentation.save(file_path)

        result = build_document_preview(file_path)

        self.assertEqual(result["kind"], "PowerPoint presentation")
        self.assertIn("Slide 1", result["content"])
        self.assertIn("Launch plan", result["content"])
        self.assertIn("Owner | Team", result["content"])

    def test_rejects_corrupted_office_document(self):
        file_path = self.path("corrupted.docx")
        with open(file_path, "wb") as output:
            output.write(b"not a zip file")

        with self.assertRaisesRegex(DocumentPreviewError, "invalid or corrupted"):
            build_document_preview(file_path)

    def test_rejects_unsafe_archive_compression_ratio(self):
        file_path = self.path("unsafe.docx")
        with zipfile.ZipFile(file_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            archive.writestr("word/document.xml", "A" * 1_000_000)

        with self.assertRaisesRegex(DocumentPreviewError, "unsafe compression ratio"):
            build_document_preview(file_path)

    def test_rejects_oversized_file_before_parsing(self):
        file_path = self.path("large.pdf")
        with open(file_path, "wb") as output:
            output.write(b"%PDF")

        with patch("core.document_preview.MAX_FILE_BYTES", 3):
            with self.assertRaisesRegex(DocumentPreviewError, "too large"):
                build_document_preview(file_path)

    def test_output_is_limited_and_marked_truncated(self):
        from docx import Document

        file_path = self.path("long.docx")
        document = Document()
        document.add_paragraph("A" * 25_000)
        document.save(file_path)

        result = build_document_preview(file_path)

        self.assertEqual(len(result["content"]), 20_000)
        self.assertTrue(result["truncated"])

    def test_rejects_unsupported_extension(self):
        file_path = self.path("sample.txt")
        with open(file_path, "w", encoding="utf-8") as output:
            output.write("hello")

        with self.assertRaisesRegex(DocumentPreviewError, "Unsupported document type"):
            build_document_preview(file_path)


if __name__ == "__main__":
    unittest.main()
