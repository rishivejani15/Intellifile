from pypdf import PdfReader
from docx import Document
from pptx import Presentation

def extract_text(path):
    try:
        if path.endswith(".pdf"):
            reader = PdfReader(path)
            return " ".join(page.extract_text() or "" for page in reader.pages)
        elif path.endswith(".docx"):
            doc = Document(path)
            return " ".join(paragraph.text for paragraph in doc.paragraphs)
        elif path.endswith(".pptx"):
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
            return text
        elif path.endswith(".txt"):
            with open(path, "r",encoding="utf-8",errors="ignore") as f:
                return f.read()
    except Exception:
        return ""